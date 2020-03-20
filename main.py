import sys
import fileinput
from pptx import Presentation

presentFile = sys.argv[1]
givenData = []
presentationData = {
    "Name": "Undefined",
    "Author": "Undefined"
}

for line in fileinput.input(presentFile):
    givenData.append(line.splitlines()[0])

presentationData["Name"] = givenData[1]
presentationData["Author"] = givenData[0]

prs = Presentation()
title_slide_layout = prs.slide_layouts[0]
slide = prs.slides.add_slide(title_slide_layout)
title = slide.shapes.title
subtitle = slide.placeholders[1]

title.text = presentationData["Name"]
subtitle.text = presentationData["Author"]

prs.save("example/test.pptx")