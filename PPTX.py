from pptx import Presentation

def get(prsData):
    prs = Presentation()
    title_slide_layout = prs.slide_layouts[0]
    slide = prs.slides.add_slide(title_slide_layout)

    slide.shapes.title.text = prsData["Title"]
    slide.placeholders[1].text = prsData["Author"]

    return prs