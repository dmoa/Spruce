import os

from PIL import Image
from pptx import Presentation
from pptx.util import Inches

EMUS_PER_INCH = 914400


def get(prsData):
    prs = Presentation()

    # layout for intro and conclusion slides
    titleTemplate = prs.slide_layouts[0]
    # layout for body slides
    bodyTemplate = prs.slide_layouts[5]

    slide = prs.slides.add_slide(titleTemplate)
    slide.shapes.title.text = prsData["Title"]
    slide.placeholders[1].text = prsData["Author"]

    for slideData in prsData["Body"]:

        slide = prs.slides.add_slide(bodyTemplate)
        slide.shapes.title.text = slideData["Title"]

        left = top = width = height = Inches(1)
        top = Inches(2)
        body = slide.shapes.add_textbox(left, top, width, height)

        body.text = slideData["Body"]

        if slideData["Image"] != "":
            im = Image.open(slideData["Image"])
            imgWidth, imgHeight = im.size
            im.close()

            _width = Inches(5) if imgWidth > imgHeight else Inches(3)
            left = Inches(9.5 - _width / EMUS_PER_INCH)

            slide.shapes.add_picture(slideData["Image"], left, top,
                                     width=_width)

            os.remove(slideData["Image"])

    return prs
